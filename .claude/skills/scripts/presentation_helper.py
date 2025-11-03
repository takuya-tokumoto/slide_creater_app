#!/usr/bin/env python3
"""
京都プレゼンテーションテンプレートを使用してスライドを作成するヘルパー関数
"""

from pptx import Presentation
from pathlib import Path


def create_presentation_from_template(template_path: str) -> Presentation:
    """
    テンプレートから新しいプレゼンテーションを作成
    
    Args:
        template_path: テンプレートファイルのパス
        
    Returns:
        Presentation: 新しいプレゼンテーションオブジェクト
    """
    return Presentation(template_path)


def get_layout_by_name(prs: Presentation, layout_name: str):
    """
    レイアウト名からレイアウトオブジェクトを取得
    
    Args:
        prs: プレゼンテーションオブジェクト
        layout_name: レイアウト名（部分一致も可）
        
    Returns:
        レイアウトオブジェクト、見つからない場合はNone
    """
    for layout in prs.slide_layouts:
        if layout_name.lower() in layout.name.lower():
            return layout
    return None


def get_layout_by_index(prs: Presentation, index: int):
    """
    インデックスからレイアウトオブジェクトを取得
    
    Args:
        prs: プレゼンテーションオブジェクト
        index: レイアウトのインデックス（0-26）
        
    Returns:
        レイアウトオブジェクト
    """
    return prs.slide_layouts[index]


def add_slide_with_layout(prs: Presentation, layout_name_or_index):
    """
    指定されたレイアウトで新しいスライドを追加
    
    Args:
        prs: プレゼンテーションオブジェクト
        layout_name_or_index: レイアウト名（文字列）またはインデックス（整数）
        
    Returns:
        追加されたスライドオブジェクト
    """
    if isinstance(layout_name_or_index, int):
        layout = get_layout_by_index(prs, layout_name_or_index)
    else:
        layout = get_layout_by_name(prs, layout_name_or_index)
    
    if layout is None:
        raise ValueError(f"レイアウトが見つかりません: {layout_name_or_index}")
    
    return prs.slides.add_slide(layout)


def list_available_layouts(prs: Presentation):
    """
    利用可能なレイアウトの一覧を表示
    
    Args:
        prs: プレゼンテーションオブジェクト
    """
    print("利用可能なレイアウト:")
    for i, layout in enumerate(prs.slide_layouts):
        print(f"  [{i}] {layout.name}")


# レイアウトインデックスの定数
LAYOUT_TITLE = 0  # タイトルまたは画面切り替え
LAYOUT_SECTION_NUMBER = 1  # 数字のセクション区切り線
LAYOUT_SECTION_DONUT = 2  # ドーナツグラフの基本セクションの区切り線
LAYOUT_QUOTE_DARK_LEFT = 3  # クォート ダーク オプション フラッシュ左揃え
LAYOUT_QUOTE_DARK_CENTER = 4  # クォート ダーク オプション中央揃え
LAYOUT_QUOTE_LIGHT_CENTER = 6  # クォート ライト オプション中央揃え
LAYOUT_BULLETS_DARK = 7  # 小さな箇条書きの付いたダークな吹き出し
LAYOUT_BULLETS_WITH_NUMBER = 8  # 小さい箇条書きの付いた吹き出し（スライド番号付き）
LAYOUT_BULLETS_7 = 9  # 小さい箇条書きの付いた吹き出し（7項目）
LAYOUT_KEYPOINTS_5 = 10  # タイトル、5つの小さな要点のテキスト
LAYOUT_50_50_RIGHT_PHOTO = 12  # 50-50 右の写真のレイアウト
LAYOUT_50_50_LEFT_PHOTO = 15  # 50-50 左フォトレイアウト
LAYOUT_BULLETS_WITH_IMAGE = 20  # 小さい箇条書きの付いた吹き出し（画像付き）
LAYOUT_BLANK = 21  # 空白
LAYOUT_BLANK_WITH_IMAGE = 22  # 画像のある空白
LAYOUT_BLANK_MINIMAL = 23  # 空白（最小限）
LAYOUT_STEPS = 24  # 手順


if __name__ == "__main__":
    # 使用例
    import sys
    
    if len(sys.argv) < 2:
        print("使用方法: python create_presentation.py <テンプレートパス>")
        sys.exit(1)
    
    template_path = sys.argv[1]
    prs = create_presentation_from_template(template_path)
    list_available_layouts(prs)
