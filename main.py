"""
スライド作成アプリ - FastAPI Backend
ES入力 → AI構成案生成 → チャット編集 → PPTX出力
"""
from fastapi import FastAPI, HTTPException
from fastapi.staticfiles import StaticFiles
from fastapi.responses import FileResponse, JSONResponse
from pydantic import BaseModel
from typing import List, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
import os
import uuid
from datetime import datetime
from anthropic import Anthropic
import json
from dotenv import load_dotenv

# .envファイルを読み込み
load_dotenv()

app = FastAPI(title="Slide Creator App")

# 静的ファイル配信
app.mount("/static", StaticFiles(directory="static"), name="static")

# エクスポートディレクトリ
EXPORT_DIR = "exports"
os.makedirs(EXPORT_DIR, exist_ok=True)

# Anthropic クライアント初期化
anthropic_client = None
try:
    api_key = os.getenv("ANTHROPIC_API_KEY")
    if api_key and api_key != "your_api_key_here":
        anthropic_client = Anthropic(api_key=api_key)
except Exception as e:
    print(f"Warning: Anthropic client initialization failed: {e}")
    print("Falling back to rule-based slide generation")


# データモデル
class Section(BaseModel):
    """ES入力のセクション"""
    title: str
    content: str


class Slide(BaseModel):
    """スライド構成"""
    title: str
    bullets: List[str]


class SlidesState(BaseModel):
    """スライド群の状態"""
    slides: List[Slide]


class GenerateRequest(BaseModel):
    """構成案生成リクエスト"""
    sections: List[Section]


class PatchRequest(BaseModel):
    """差分編集リクエスト"""
    slides: List[Slide]
    prompt: str


class ExportRequest(BaseModel):
    """PPTX出力リクエスト"""
    slides: List[Slide]


# ルートエンドポイント
@app.get("/")
async def root():
    """ES入力フォームページを返す"""
    return FileResponse("static/index.html")


class MessageLineSlide(BaseModel):
    """メッセージラインとタイトルのみを含むスライド"""
    title: str
    message_line: str  # 40文字以内の核心メッセージ


async def generate_message_lines(sections: List[Section]) -> List[MessageLineSlide]:
    """
    Step1: メッセージラインを生成する
    各スライドの核心メッセージ（1~2行・80文字以内）とタイトルを決定
    """
    sections_text = "\n\n".join([
        f"【{section.title}】\n{section.content}"
        for section in sections
    ])

    prompt = f"""以下の自己PR・ES情報を分析し、効果的なプレゼンテーションスライドのメッセージラインを決定してください。

# 入力情報
{sections_text}

# 要件
1. 入力情報全体を分析し、各スライドの核心メッセージを1~2行で定義（80文字以内）
2. 事実と示唆を統合した形で記述（プレフィックス不要）
3. 最初のスライドはタイトルスライド（全体の目的を明示）
4. 最後にまとめスライド（結論と次アクション）
5. 全体で5-8枚程度のスライド
6. 情報を適切にグループ化し、ストーリー性を持たせる

# 出力形式
以下のJSON配列形式で返してください：

```json
{{
  "slides": [
    {{
      "title": "スライドのタイトル",
      "message_line": "核心メッセージ（80文字以内）"
    }}
  ]
}}
```

JSON配列のみを返してください（説明文は不要）。"""

    try:
        response = anthropic_client.messages.create(
            model="claude-sonnet-4-20250514",
            max_tokens=2000,
            messages=[
                {"role": "user", "content": prompt}
            ]
        )

        # レスポンスからJSONを抽出
        content = response.content[0].text

        if "```json" in content:
            content = content.split("```json")[1].split("```")[0].strip()
        elif "```" in content:
            content = content.split("```")[1].split("```")[0].strip()

        # JSONをパース
        data = json.loads(content)
        message_lines = [MessageLineSlide(**slide) for slide in data.get("slides", [])]
        return message_lines

    except Exception as e:
        print(f"Message line generation error: {e}")
        raise HTTPException(
            status_code=500,
            detail=f"メッセージライン生成中にエラーが発生しました: {str(e)}"
        )


async def generate_slide_bodies_with_skills(
    message_line_slides: List[MessageLineSlide],
    sections: List[Section]
) -> List[Slide]:
    """
    Step2: Claude Skills（tools）を使用してボディ部分を生成
    メッセージラインと元データから、スライドテンプレートに従ってボディを抽出・生成
    """
    sections_text = "\n\n".join([
        f"【{section.title}】\n{section.content}"
        for section in sections
    ])

    # Skills（tools）の定義：ボディ生成用の構造化出力スキーマ
    slide_body_tool = {
        "name": "generate_slide_body",
        "description": "メッセージラインに基づいて、元データからボディ部分を抽出・生成する",
        "input_schema": {
            "type": "object",
            "properties": {
                "bullets": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "メッセージラインを裏付ける箇条書き（3-5個）。元データから具体的な根拠・事例・データを引用"
                }
            },
            "required": ["bullets"]
        }
    }

    all_slides = []

    # 各メッセージラインに対してボディを生成
    for msg_slide in message_line_slides:
        prompt = f"""以下の情報を基に、スライドのボディ部分を生成してください。

# 元データ（入力情報）
{sections_text}

# このスライドの情報
- タイトル: {msg_slide.title}
- メッセージライン: {msg_slide.message_line}

# 要件
1. メッセージライン「{msg_slide.message_line}」を裏付ける情報を元データから抽出
2. 元データに記載されている具体的な根拠・事例・データを優先的に使用
3. 3-5個の箇条書きで構成（各項目は30文字以内が目安）
4. 元データにない情報は極力避け、入力情報に忠実に基づく
5. 以下の順序で構成：
   - 1つ目: 元データからの根拠・背景情報
   - 2つ目: 元データからの具体例・データ
   - 3つ目: メッセージラインを補足する詳細説明
   - 4つ目以降: （必要に応じて）行動項目や検討ポイント

generate_slide_bodyツールを使用してボディを生成してください。"""

        try:
            response = anthropic_client.messages.create(
                model="claude-sonnet-4-20250514",
                max_tokens=2000,
                messages=[
                    {"role": "user", "content": prompt}
                ],
                tools=[slide_body_tool],
                tool_choice={"type": "tool", "name": "generate_slide_body"}
            )

            # Toolsの結果を取得
            if response.content and len(response.content) > 0:
                tool_result = response.content[0]
                if hasattr(tool_result, 'type') and tool_result.type == "tool_use":
                    if tool_result.name == "generate_slide_body":
                        # inputはdictまたは適切な形式で提供される
                        if isinstance(tool_result.input, dict):
                            body_bullets = tool_result.input.get("bullets", [])
                        else:
                            # フォールバック: inputがdictでない場合
                            body_bullets = getattr(tool_result.input, "bullets", [])
                        
                        # メッセージラインを先頭に追加
                        full_bullets = [msg_slide.message_line] + body_bullets
                        
                        all_slides.append(Slide(
                            title=msg_slide.title,
                            bullets=full_bullets
                        ))
                    else:
                        # 想定外のtool名
                        print(f"Warning: Unexpected tool name '{tool_result.name}'")
                        all_slides.append(Slide(
                            title=msg_slide.title,
                            bullets=[msg_slide.message_line, "詳細情報を元データから抽出してください"]
                        ))
            else:
                # Toolsが使えない場合のフォールバック
                print(f"Warning: Tools not used for slide '{msg_slide.title}', using fallback")
                all_slides.append(Slide(
                    title=msg_slide.title,
                    bullets=[msg_slide.message_line, "詳細情報を元データから抽出してください"]
                ))

        except Exception as e:
            print(f"Body generation error for slide '{msg_slide.title}': {e}")
            # エラー時はメッセージラインのみでスライドを作成
            all_slides.append(Slide(
                title=msg_slide.title,
                bullets=[msg_slide.message_line]
            ))

    return all_slides


async def generate_slides_with_llm(sections: List[Section]) -> List[Slide]:
    """
    LLMを使用してスライド構成を生成（2段階アプローチ）
    Step1: メッセージライン生成
    Step2: Skills機能でボディ生成
    """
    # Step1: メッセージラインを生成
    message_line_slides = await generate_message_lines(sections)
    
    # Step2: メッセージラインと元データからボディを生成
    slides = await generate_slide_bodies_with_skills(message_line_slides, sections)
    
    return slides


@app.post("/generate")
async def generate_slides(request: GenerateRequest) -> SlidesState:
    """
    ES入力から構成案を生成（LLMのみ使用）
    """
    if not anthropic_client:
        raise HTTPException(
            status_code=500,
            detail="Anthropic API キーが設定されていません。.envファイルにANTHROPIC_API_KEYを設定してください。"
        )

    slides = await generate_slides_with_llm(request.sections)
    return SlidesState(slides=slides)


@app.post("/patch")
async def patch_slides(request: PatchRequest) -> SlidesState:
    """
    チャット入力でスライドを編集
    簡易実装：プロンプトに応じた操作を解析
    """
    slides = request.slides.copy()
    prompt = request.prompt.lower()

    # プロンプト解析（簡易版）
    if "削除" in prompt or "消して" in prompt or "delete" in prompt:
        # 最後のスライドを削除（タイトル以外）
        if len(slides) > 1:
            slides.pop()

    elif "追加" in prompt or "add" in prompt:
        # 新しいスライドを追加
        slides.append(Slide(
            title="新しいスライド",
            bullets=["内容を編集してください"]
        ))

    elif "タイトル" in prompt and "変更" in prompt:
        # 最初のスライドのタイトルを変更
        if slides and "→" in prompt:
            new_title = prompt.split("→")[1].strip()
            slides[0].title = new_title

    elif "箇条書き" in prompt or "内容" in prompt:
        # 箇条書きを追加
        if len(slides) > 1:
            new_bullet = prompt.replace("箇条書き", "").replace("追加", "").strip()
            if new_bullet:
                slides[-1].bullets.append(new_bullet)

    else:
        # デフォルト：最後のスライドに内容を追加
        if slides:
            slides[-1].bullets.append(f"💡 {request.prompt}")

    return SlidesState(slides=slides)


@app.post("/export")
async def export_pptx(request: ExportRequest) -> dict:
    """
    PPTXファイルを生成してダウンロードURLを返す
    """
    # プレゼンテーション作成
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    for slide_data in request.slides:
        # タイトルと内容のレイアウト
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)

        # タイトル設定
        title = slide.shapes.title
        title.text = slide_data.title

        # 箇条書き設定
        if slide_data.bullets:
            body = slide.placeholders[1]
            text_frame = body.text_frame
            text_frame.clear()

            for i, bullet in enumerate(slide_data.bullets):
                if i == 0:
                    text_frame.text = bullet
                    run = text_frame.paragraphs[0].runs[0]
                    run.font.bold = True  # 1行目（メッセージライン）を太字
                else:
                    p = text_frame.add_paragraph()
                    p.text = bullet
                    p.level = 0

    # ファイル保存
    filename = f"slide_{uuid.uuid4().hex[:8]}.pptx"
    filepath = os.path.join(EXPORT_DIR, filename)
    prs.save(filepath)

    return {
        "download_url": f"/download/{filename}",
        "filename": filename
    }


@app.get("/download/{filename}")
async def download_file(filename: str):
    """
    生成されたPPTXファイルをダウンロード
    """
    filepath = os.path.join(EXPORT_DIR, filename)

    if not os.path.exists(filepath):
        raise HTTPException(status_code=404, detail="ファイルが見つかりません")

    return FileResponse(
        filepath,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=filename
    )


if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="127.0.0.1", port=8000)
